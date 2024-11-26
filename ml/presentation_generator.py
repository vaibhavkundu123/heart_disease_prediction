from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import logging
from typing import List, Dict, Any

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PresentationGenerator:
    def __init__(self):
        self.prs = Presentation()
        # Define constants
        self.TITLE_COLOR = RGBColor(7, 66, 154)  # #07429a
        self.SUBTITLE_COLOR = RGBColor(166, 173, 22)  # #a6ad16
        self.TEXT_COLOR = RGBColor(51, 51, 51)
        self.ACCENT_COLOR = RGBColor(76, 175, 80)  # #4CAF50
        
    def create_title_slide(self) -> None:
        """Creates the title slide"""
        title_slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        title.text = "Patient Health Disease Prediction System"
        subtitle.text = "A Web-Based Healthcare Decision Support System"
        
        self._format_text_frame(title.text_frame, self.TITLE_COLOR, size=44)
        self._format_text_frame(subtitle.text_frame, self.SUBTITLE_COLOR, size=28)

    def create_content_slide(self, title: str, content: List[str]) -> None:
        """Creates a content slide with bullet points"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = title
        self._format_text_frame(title_shape.text_frame, self.TITLE_COLOR, size=36)
        
        tf = body_shape.text_frame
        for line in content:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith('###'):
                self._format_paragraph(p, self.ACCENT_COLOR, size=20, bold=True, level=0)
                p.text = line.replace('###', '').strip()
            elif line.startswith('-'):
                self._format_paragraph(p, self.TEXT_COLOR, size=18, level=1)
            else:
                self._format_paragraph(p, self.TEXT_COLOR, size=18, bold=True, level=0)

    def _format_text_frame(self, text_frame, color, size: int = 18) -> None:
        """Formats text frame with given properties"""
        for paragraph in text_frame.paragraphs:
            paragraph.font.color.rgb = color
            paragraph.font.size = Pt(size)

    def _format_paragraph(self, paragraph, color, size: int = 18, bold: bool = False, level: int = 0) -> None:
        """Formats paragraph with given properties"""
        paragraph.font.color.rgb = color
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.level = level

    def generate_presentation(self) -> None:
        """Generates complete presentation"""
        try:
            logger.info("Starting presentation generation...")
            
            # Create title slide
            self.create_title_slide()
            
            # Define all slides content
            slides_content = {
                "Project Overview": [
                    "### Key Components",
                    "- Web interface for data input and results visualization",
                    "- Machine learning model for multi-disease prediction",
                    "- Secure user authentication and data management",
                    "- Historical tracking of predictions",
                    "### Core Features", 
                    "- Multi-disease classification using ML algorithms",
                    "- Probability-based confidence scoring",
                    "- Comprehensive patient data collection",
                    "- Historical prediction tracking"
                ],
                "Technology Stack": [
                    "### Frontend Technologies",
                    "- HTML5 for document structure",
                    "- CSS3 for styling and responsive design", 
                    "- JavaScript for client-side interactivity",
                    "- Bootstrap 5 for responsive UI components",
                    "### Backend Technologies",
                    "- Django 4.0 web framework",
                    "- Python 3.9+ for application logic",
                    "- SQLite database for data persistence",
                    "### Machine Learning Stack",
                    "- scikit-learn for ML algorithms",
                    "- pandas for data manipulation",
                    "- numpy for numerical operations",
                    "- matplotlib/seaborn for visualization"
                ],
                "System Architecture": [
                    "### Django MVT Architecture",
                    "- Models: PredictionResult, User, PatientData",
                    "- Views: Authentication, Prediction, History",
                    "- Templates: Forms, Results, Dashboard",
                    "### Key Components",
                    "- Authentication System",
                    "- Prediction Engine", 
                    "- Data Processing Pipeline",
                    "- Results Storage",
                    "- API Layer"
                ],
                "Key Features": [
                    "### User Management",
                    "- Secure registration and login",
                    "- Password reset capability",
                    "- Role-based access control",
                    "- Profile management",
                    "### Prediction System",
                    "- Real-time disease prediction",
                    "- Multiple disease detection", 
                    "- Confidence scoring",
                    "- Recommendation generation"
                ],
                "Data Collection": [
                    "### Patient Metrics",
                    "- Height (cm)",
                    "- Weight (kg)", 
                    "- Temperature (°C)",
                    "- Heart Rate (bpm)",
                    "- Blood Pressure",
                    "- Cholesterol (mg/dL)",
                    "- Blood Sugar (mg/dL)",
                    "### Medical Information",
                    "- Current Symptoms",
                    "- Existing Conditions",
                    "- Family History",
                    "- Laboratory Results",
                    "- Lifestyle Factors"
                ],
                "Data Processing": [
                    "### Preprocessing Pipeline",
                    "- Missing value imputation",
                    "- Feature scaling with StandardScaler",
                    "- Categorical encoding",
                    "- Blood pressure splitting",
                    "- Feature normalization",
                    "### Data Cleaning",
                    "- Outlier detection",
                    "- Data type conversion",
                    "- Date/time standardization",
                    "- Invalid value handling"
                ],
                "Machine Learning Model": [
                    "### Algorithm Details",
                    "- Multi-label K-Nearest Neighbors",
                    "- Correlation-based feature learning",
                    "- Cross-validation k=5 folds",
                    "- Feature importance ranking",
                    "### Model Metrics",
                    "- Accuracy: ~85-90%",
                    "- Precision: 87%",
                    "- Recall: 86%",
                    "- F1-Score: 86.5%"
                ],
                "User Interface - Home": [
                    "### Design Elements",
                    "- Clean medical interface theme",
                    "- Blue (#07429a) primary color scheme",
                    "- Responsive Bootstrap layout",
                    "- Quick access navigation",
                    "### Components",
                    "- Login/Register buttons",
                    "- Welcome message",
                    "- System description",
                    "- Help information"
                ],
                "User Interface - Prediction Form": [
                    "### Form Fields",
                    "- Patient Demographics",
                    "- Vital Signs (Height/Weight/Temperature)",
                    "- Blood Pressure",
                    "- Heart Rate",
                    "### Medical Data",
                    "- Symptoms",
                    "- Existing Conditions",
                    "- Lab Results",
                    "- Family History"
                ],
                "User Interface - Results": [
                    "### Results Display",
                    "- Primary predicted disease",
                    "- Confidence percentage",
                    "- Alternative predictions",
                    "- Supporting evidence",
                    "### Patient Data Summary",
                    "- Input parameters",
                    "- Key metrics",
                    "- Risk factors",
                    "- Recommendations"
                ],
                "User Interface - History": [
                    "### History Features",
                    "- Date/time of prediction",
                    "- Predicted disease",
                    "- Confidence score",
                    "- Patient metrics",
                    "### Functionality",
                    "- Sortable columns",
                    "- Filtering options",
                    "- Export capability",
                    "- Trend analysis"
                ],
                "Security Features": [
                    "### Authentication",
                    "- Django user authentication",
                    "- Password hashing",
                    "- Session management",
                    "- Login rate limiting",
                    "### Data Protection",
                    "- CSRF protection",
                    "- XSS prevention",
                    "- SQL injection protection",
                    "- Encrypted storage"
                ],
                "Model Training": [
                    "### Training Process",
                    "- Data splitting (80/20)",
                    "- Feature scaling",
                    "- Cross-validation",
                    "- Hyperparameter tuning",
                    "### Feature Importance",
                    "- Heart Rate",
                    "- Blood Pressure",
                    "- Cholesterol",
                    "- Symptoms",
                    "- Family History"
                ],
                "Performance Metrics": [
                    "### Model Evaluation",
                    "- Training accuracy: 88%",
                    "- Validation accuracy: 85%",
                    "- Test accuracy: 86%",
                    "- AUC-ROC: 0.89",
                    "### Cross-validation",
                    "- 5-fold CV scores",
                    "- Mean accuracy: 0.87",
                    "- Standard deviation: 0.02"
                ],
                "Testing Strategy": [
                    "### Test Coverage",
                    "- Unit tests for models",
                    "- Integration tests for views",
                    "- Form validation tests",
                    "- API endpoint tests",
                    "### Quality Assurance",
                    "- Automated testing",
                    "- Code reviews",
                    "- Performance monitoring",
                    "- Security scanning"
                ],
                "Deployment": [
                    "### Requirements",
                    "- Python 3.9+",
                    "- Django 4.0+",
                    "- SQLite/PostgreSQL",
                    "- Web server (Nginx/Apache)",
                    "### Configuration",
                    "- Environment setup",
                    "- Database initialization",
                    "- Static files",
                    "- Security settings"
                ],
                "Future Enhancements": [
                    "### Planned Features",
                    "- Mobile responsive design",
                    "- REST API endpoints",
                    "- Additional ML models",
                    "- Enhanced visualizations",
                    "### Integration Options",
                    "- EMR systems",
                    "- Lab systems",
                    "- Pharmacy systems",
                    "- Billing systems"
                ],
                "Benefits": [
                    "### Clinical Impact",
                    "- Faster diagnosis",
                    "- Evidence-based decisions",
                    "- Multi-disease detection",
                    "- Risk assessment",
                    "### Operational Benefits",
                    "- Efficient workflow",
                    "- Digital records",
                    "- Historical tracking",
                    "- Decision support"
                ],
                "Conclusion": [
                    "### Key Achievements",
                    "- Multi-disease prediction",
                    "- User-friendly interface",
                    "- Secure data handling",
                    "- Historical tracking",
                    "### Future Scope",
                    "- Model improvements",
                    "- Feature expansion",
                    "- Mobile platform",
                    "- System integration"
                ]
            }
            
            # Generate content slides
            for title, content in slides_content.items():
                logger.info(f"Creating slide: {title}")
                self.create_content_slide(title, content)
            
            # Save presentation
            output_file = "patient_health_prediction.pptx"
            self.prs.save(output_file)
            logger.info(f"Presentation saved successfully: {output_file}")
            
        except Exception as e:
            logger.error(f"Error generating presentation: {str(e)}")
            raise

def main():
    """Main execution function"""
    try:
        generator = PresentationGenerator()
        generator.generate_presentation()
    except Exception as e:
        logger.error(f"Failed to generate presentation: {str(e)}")
        raise

if __name__ == "__main__":
    main()